#!/usr/bin/env python
# coding: utf-8

# In[ ]:


import openai
from pptx import Presentation
from pptx.util import Inches
import os
import time

# Function to collect user inputs
def get_user_inputs():
    api_key = input("Enter your OpenAI API key: ")
    subject = input("Enter the subject of the presentation: ")
    audience = input("Specify the audience: ")
    num_images = int(input("Enter the number of images: "))
    num_slides = int(input("Enter the number of slides: "))
    return api_key, subject, audience, num_images, num_slides


# In[ ]:


# Creating the prompt based on user inputs
def create_prompt(subject, audience, num_images, num_slides):
    prompt = f"""
    You're skilled in presentation creation. You create highly effective and striking PowerPoint Presentations. Your task is to create a short presentation about {subject} aimed at a {audience} audience.
    
    You will do this step-by-step:
    1. Retrieve the data by searching the web.
    2. Write clear and concise text content for {num_slides} slides: a title page (engaging short title, 2-3 words, and sub-title, 4-5 words), introduction page (presentation of the context), description of the data page (presenting the data), trend observation page (trends that are observable from the data), conclusion page. Ensure thoughtful formatting with spacing, lists and short lines.
    3. Generate {num_images} supporting images, each supporting the slides. Use landscape mode and a minimalist illustration style.
    4. Let me review the images and text before proceeding.
    5. Generate and execute code to convert the images to PNG format.
    6. Create the PowerPoint presentation with {num_slides} slides that include the generated text and images.
    """
    return prompt


# In[ ]:


# Function to generate content using the OpenAI API with error handling
def generate_presentation_content(prompt):
    max_retries = 5
    retries = 0
    while retries < max_retries:
        try:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant.You will generate content to the point and ethically."},
                    {"role": "user", "content": prompt}
                ]
            )
            return response.choices[0].message['content'].strip()
        except openai.error.RateLimitError:
            retries += 1
            print(f"Rate limit exceeded. Retrying in {2 ** retries} seconds...")
            time.sleep(2 ** retries)
    raise Exception("Rate limit exceeded. Please try again later.")


# In[ ]:


# Function to create a PowerPoint presentation
def create_presentation(content, num_slides):
    prs = Presentation()
    
    # Split the content into slides
    slides_content = content.split('\n\n')[:num_slides]
    
    # Add slides
    for i, slide_content in enumerate(slides_content):
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = f"Slide {i+1}"
        body.text = slide_content
    
    # Save the presentation
    prs.save('presentation.pptx')
    print("PowerPoint presentation created as 'presentation.pptx'.")


# In[ ]:


# Main function to execute the steps
def main():
    api_key, subject, audience, num_images, num_slides = get_user_inputs()
    openai.api_key = api_key
    
    prompt = create_prompt(subject, audience, num_images, num_slides)
    
    content = generate_presentation_content(prompt)
    print("Generated Content:\n", content)
    
    # Assume user reviews and confirms the content and images here
    input("Review the generated content above and press Enter to continue...")
    
    create_presentation(content, num_slides)

if __name__ == "__main__":
    main()


# In[ ]:




